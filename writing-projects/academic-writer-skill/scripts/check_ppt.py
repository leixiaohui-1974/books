#!/usr/bin/env python3
"""
演示文稿(PPT)质量自动检查脚本
用法: python3 check_ppt.py <文件路径>
支持: .pptx (python-pptx) 和 .md (Markdown大纲) 两种输入

配合 writing-projects/academic-writer-skill/SKILL.md §7.7 使用
风格指南: references/ppt_style_guide.md
"""
import re
import sys
import json
import argparse
from pathlib import Path

# ============================================================
# 颜色/字体规范常量 (来自PPT风格指南)
# ============================================================
APPROVED_COLORS = {
    # 主色-科技蓝
    '#002060', '#0F62AD', '#2E75B5',
    # 强调色-中国红
    '#C00000', '#FF0000',
    # 辅助色
    '#FFFF00', '#7030A0',
    # 正文色
    '#000000', '#3F3F3F',
    # 背景色
    '#FFFFFF',
}

APPROVED_CN_FONTS = {'微软雅黑', 'Microsoft YaHei', 'Microsoft YaHei UI'}
APPROVED_EN_FONTS = {'Times New Roman', 'Arial', 'Calibri'}

FONT_SIZE_RULES = {
    'cover_title': (36, 44),
    'section_title': (32, 40),
    'content_title': (24, 28),
    'body': (16, 20),
    'label': (12, 16),
    'en_number': (12, 18),
}

# ============================================================
# PPTX 检查 (需要python-pptx)
# ============================================================

def check_pptx(filepath):
    """检查 .pptx 文件"""
    try:
        from pptx import Presentation
        from pptx.util import Pt, Emu
        from pptx.dml.color import RGBColor
    except ImportError:
        print("⚠️  python-pptx未安装，使用: pip install python-pptx --break-system-packages")
        return None, None
    
    prs = Presentation(filepath)
    issues = []
    stats = {
        '总页数': len(prs.slides),
        '幻灯片宽度': f'{prs.slide_width / 914400:.1f}"',
        '幻灯片高度': f'{prs.slide_height / 914400:.1f}"',
    }
    
    # 1. 检查比例 (16:9 = 13.333" x 7.5")
    width_inches = prs.slide_width / 914400
    height_inches = prs.slide_height / 914400
    ratio = width_inches / height_inches
    if abs(ratio - 16/9) > 0.1:
        issues.append({
            'level': 'SEVERE', 'category': '页面比例',
            'message': f'页面比例为 {ratio:.2f}:1，非16:9宽屏({width_inches:.1f}×{height_inches:.1f}英寸)',
            'suggestion': '设置为16:9宽屏(33.867cm × 19.05cm)'
        })
    
    # 2. 逐页检查
    color_violations = []
    font_violations = []
    no_notes_pages = []
    text_heavy_pages = []
    
    for idx, slide in enumerate(prs.slides, 1):
        # 2a. 检查Speaker Notes
        if not slide.has_notes_slide or not slide.notes_slide.notes_text_frame.text.strip():
            no_notes_pages.append(idx)
        
        # 2b. 检查颜色和字体
        slide_text_total = 0
        slide_has_image = False
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        slide_text_total += len(run.text)
                        # 检查字体颜色
                        if run.font.color and run.font.color.rgb:
                            color_hex = f'#{run.font.color.rgb}'
                            if color_hex.upper() not in {c.upper() for c in APPROVED_COLORS}:
                                color_violations.append((idx, color_hex, run.text[:20]))
                        # 检查字体
                        if run.font.name:
                            fname = run.font.name
                            if fname not in APPROVED_CN_FONTS and fname not in APPROVED_EN_FONTS:
                                font_violations.append((idx, fname, run.text[:20]))
            
            # 检查图片
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                slide_has_image = True
        
        # 2c. 检查文字墙 (超过200字无图片)
        if slide_text_total > 200 and not slide_has_image:
            text_heavy_pages.append((idx, slide_text_total))
    
    # 汇总问题
    if no_notes_pages:
        if len(no_notes_pages) > 3:
            issues.append({
                'level': 'SEVERE', 'category': 'Speaker Notes',
                'message': f'{len(no_notes_pages)}页缺少讲稿: 第{no_notes_pages[:5]}...页',
                'suggestion': '每页须有Speaker Notes（含转场语），覆盖率须100%'
            })
        else:
            issues.append({
                'level': 'WARNING', 'category': 'Speaker Notes',
                'message': f'第{no_notes_pages}页缺少讲稿',
                'suggestion': '补充Speaker Notes'
            })
    
    stats['Notes覆盖率'] = f"{(len(prs.slides)-len(no_notes_pages))/len(prs.slides)*100:.0f}%"
    
    # 去重颜色违规
    unique_bad_colors = set(c for _, c, _ in color_violations)
    if unique_bad_colors:
        issues.append({
            'level': 'WARNING', 'category': '配色合规',
            'message': f'发现{len(unique_bad_colors)}种非规范颜色: {unique_bad_colors}',
            'suggestion': '仅使用科技蓝/中国红/正文黑灰体系'
        })
    
    # 去重字体违规
    unique_bad_fonts = set(f for _, f, _ in font_violations)
    if unique_bad_fonts:
        issues.append({
            'level': 'WARNING', 'category': '字体合规',
            'message': f'发现非规范字体: {unique_bad_fonts}',
            'suggestion': '中文用微软雅黑，英文用Times New Roman/Arial'
        })
    
    if text_heavy_pages:
        pages_str = ', '.join(f'第{p}页({c}字)' for p, c in text_heavy_pages[:5])
        issues.append({
            'level': 'WARNING', 'category': '信息密度',
            'message': f'文字墙页面: {pages_str}',
            'suggestion': '图文比应≥4:6，添加图表/图片或拆分内容'
        })
    
    return issues, stats

# ============================================================
# Markdown大纲检查 (PPT脚本生成前检查)
# ============================================================

def check_md_outline(filepath):
    """检查 Markdown 格式的PPT大纲"""
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()
    
    issues = []
    lines = content.split('\n')
    
    slides = []
    current_slide = None
    for line in lines:
        if re.match(r'^#{1,2}\s|^Slide\s*\d+|^---', line):
            if current_slide:
                slides.append(current_slide)
            current_slide = {'title': line.strip(), 'content': '', 'has_notes': False}
        elif current_slide:
            current_slide['content'] += line + '\n'
            if 'Notes:' in line or 'Speaker Notes' in line or '讲稿' in line:
                current_slide['has_notes'] = True
    if current_slide:
        slides.append(current_slide)
    
    stats = {'总页数': len(slides)}
    
    # 检查基本结构
    if len(slides) < 5:
        issues.append({
            'level': 'WARNING', 'category': '页数',
            'message': f'仅{len(slides)}页，内容可能不足',
            'suggestion': '通常至少需要8页以上'
        })
    
    # 检查封面/结尾
    if slides:
        first_title = slides[0]['title'].lower()
        if '封面' not in first_title and 'cover' not in first_title and 'slide 1' not in first_title:
            issues.append({
                'level': 'CHECK', 'category': '结构',
                'message': '未检测到明确的封面页标识',
                'suggestion': '第一页应为封面(项目名称+Logo+汇报人+日期)'
            })
        
        last_title = slides[-1]['title'].lower()
        if '谢谢' not in last_title and 'thank' not in last_title and '结尾' not in last_title:
            issues.append({
                'level': 'CHECK', 'category': '结构',
                'message': '未检测到明确的结尾页标识',
                'suggestion': '最后一页应为结尾页(谢谢/Thank You)'
            })
    
    # 检查过渡页
    transition_count = sum(1 for s in slides if '过渡' in s['title'] or '章节' in s['title'])
    if len(slides) > 15 and transition_count < 2:
        issues.append({
            'level': 'WARNING', 'category': '节奏',
            'message': f'{len(slides)}页中仅{transition_count}个过渡页',
            'suggestion': '每章节应有过渡页分隔，起到节奏引导作用'
        })
    
    # 检查连续高密度页面
    consecutive_content = 0
    for s in slides:
        if '过渡' not in s['title'] and '封面' not in s['title'] and '谢谢' not in s['title']:
            consecutive_content += 1
            if consecutive_content > 5:
                issues.append({
                    'level': 'WARNING', 'category': '节奏',
                    'message': f'连续{consecutive_content}页高密度内容无过渡页',
                    'suggestion': '连续5页以上须插入呼吸页/过渡页'
                })
                break
        else:
            consecutive_content = 0
    
    # 检查Notes覆盖率
    no_notes = [i+1 for i, s in enumerate(slides) if not s['has_notes']]
    if no_notes:
        stats['Notes覆盖率'] = f"{(len(slides)-len(no_notes))/max(len(slides),1)*100:.0f}%"
        if len(no_notes) > len(slides) * 0.3:
            issues.append({
                'level': 'SEVERE', 'category': 'Speaker Notes',
                'message': f'{len(no_notes)}/{len(slides)}页缺少讲稿',
                'suggestion': '每页须有Speaker Notes，覆盖率须100%'
            })
    else:
        stats['Notes覆盖率'] = '100%'
    
    return issues, stats

# ============================================================
# 时长检查
# ============================================================

def check_duration(slide_count, target_minutes=None):
    """检查页数与时长匹配"""
    issues = []
    if target_minutes:
        expected_min = target_minutes * 0.6  # 每页约1-1.5分钟
        expected_max = target_minutes * 1.2
        if slide_count < expected_min:
            issues.append({
                'level': 'WARNING', 'category': '时长匹配',
                'message': f'{slide_count}页对于{target_minutes}分钟汇报偏少(建议{int(expected_min)}-{int(expected_max)}页)',
                'suggestion': '增加内容页或拆分复杂页面'
            })
        elif slide_count > expected_max:
            issues.append({
                'level': 'WARNING', 'category': '时长匹配',
                'message': f'{slide_count}页对于{target_minutes}分钟汇报偏多',
                'suggestion': '合并相关页面或精简内容'
            })
    return issues

# ============================================================
# 主入口
# ============================================================

def run_checks(filepath, target_minutes=None):
    """执行全部检查"""
    path = Path(filepath)
    
    if path.suffix == '.pptx':
        issues, stats = check_pptx(filepath)
        if issues is None:
            return False
    elif path.suffix in ('.md', '.txt'):
        issues, stats = check_md_outline(filepath)
    else:
        print(f"❌ 不支持的文件格式: {path.suffix}")
        print("   支持: .pptx, .md")
        return False
    
    # 时长检查
    slide_count = stats.get('总页数', 0)
    if target_minutes:
        issues.extend(check_duration(slide_count, target_minutes))
    
    # 输出报告
    print(f"\n{'='*65}")
    print(f" 🎯 演示文稿质量检查报告")
    print(f"{'='*65}")
    
    print(f"\n📊 基本统计:")
    for k, v in stats.items():
        print(f"   {k}: {v}")
    if target_minutes:
        print(f"   目标时长: {target_minutes}分钟")
        print(f"   预估语速: 每页{target_minutes/max(slide_count,1):.1f}分钟")
    
    severe = [i for i in issues if i['level'] == 'SEVERE']
    warning = [i for i in issues if i['level'] == 'WARNING']
    check = [i for i in issues if i['level'] == 'CHECK']
    
    print(f"\n共发现 {len(issues)} 个问题 "
          f"(🔴严重:{len(severe)} 🟡警告:{len(warning)} ⚪待确认:{len(check)})")
    
    # 风格合规评分
    total_items = 20  # 设计审查总检查项
    deductions = len(severe) * 3 + len(warning) * 1
    style_score = max(0, total_items - deductions)
    compliance = f"{style_score}/{total_items}"
    print(f"   风格合规: {compliance}")
    
    if severe:
        print(f"\n🔴 严重问题（必须修正）:")
        print(f"{'—'*55}")
        for i, item in enumerate(severe, 1):
            print(f"  {i}. [{item['category']}] {item['message']}")
            print(f"     → {item['suggestion']}")
    
    if warning:
        print(f"\n🟡 警告（建议修正）:")
        print(f"{'—'*55}")
        for i, item in enumerate(warning, 1):
            print(f"  {i}. [{item['category']}] {item['message']}")
            print(f"     → {item['suggestion']}")
    
    if check:
        print(f"\n⚪ 需人工确认:")
        print(f"{'—'*55}")
        for i, item in enumerate(check, 1):
            print(f"  {i}. [{item['category']}] {item['message']}")
            print(f"     → {item['suggestion']}")
    
    print(f"\n{'='*65}")
    if severe:
        print("❌ 存在严重问题，PPT尚不能提交评审。")
    elif warning:
        print("⚠️  建议修正警告项后提交评审。")
    else:
        print("✅ PPT检查通过，可提交评审。")
    print()
    
    return len(severe) == 0

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='演示文稿质量检查')
    parser.add_argument('filepath', help='PPT文件路径(.pptx或.md大纲)')
    parser.add_argument('--duration', type=int, help='目标汇报时长(分钟)')
    args = parser.parse_args()
    
    passed = run_checks(args.filepath, args.duration)
    sys.exit(0 if passed else 1)
